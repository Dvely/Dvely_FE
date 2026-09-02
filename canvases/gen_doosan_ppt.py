#!/usr/bin/env python3
"""DT_두산에너빌리티 PBL ①~⑤ 발표용 16:9 PPTX."""

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn, nsmap
from pptx.util import Emu, Inches, Pt

OUT = "/Users/kimtaewoo/Documents/GitHub/Dvely_FE/canvases/doosan-pbl-plan-v3.pptx"
FONT = "Apple SD Gothic Neo"
MONO = "Menlo"

def C(h):
    return RGBColor((h >> 16) & 0xFF, (h >> 8) & 0xFF, h & 0xFF)

NAVY = RGBColor(0x1E, 0x4A, 0x7A)
NAVY2 = RGBColor(0x0F, 0x27, 0x44)
NAVY3 = C(0x163A62)
ORANGE = C(0xC2410C)
GOLD = C(0xB45309)
TEAL = C(0x0F766E)
GREEN = C(0x166534)
RED = C(0x9A3412)
TEXT = C(0x1A1A1A)
MUTED = C(0x4B5563)
LINE = C(0xD1D5DB)
WHITE = C(0xFFFFFF)
LIGHT = C(0xF4F7FB)
PALE = C(0xEEF3E8)
WARN_BG = C(0xFFF7ED)
TEAL_BG = C(0xECFDF5)
GOLD_BG = C(0xFFFBEB)
BLUE_BG = C(0xE8F0F8)
ROW_ALT = C(0xF8FAFC)
AMBER = C(0xFCD34D)
SKY = C(0xBFDBFE)
SKY2 = C(0x93C5FD)
SLATE = C(0xF1F5F9)
NAVY_CARD = C(0x1A3A5C)

W = Inches(13.333)
H = Inches(7.5)
SLIDES = 12


def ea_font(run, name=FONT):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", name)


def set_run(run, size, bold=False, color=TEXT, name=FONT, italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = name
    ea_font(run, name)


def add_shape(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def add_round(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.adjustments[0] = 0.12
    sh.shadow.inherit = False
    return sh


def tb(shape, text, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE, name=FONT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.anchor = valign
    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.clear()
        run = p.add_run()
        run.text = line
        set_run(run, size, bold, color, name)
    return tf


def add_text(slide, l, t, w, h, text, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, name=FONT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tb(box, text, size, bold, color, align, valign, name)
    return box


def header(slide, section, title, n):
    add_shape(slide, 0, 0, W, Inches(0.78), NAVY)
    add_shape(slide, 0, 0, Inches(0.12), Inches(0.78), GOLD)
    add_text(slide, Inches(0.32), Inches(0.08), Inches(10.4), Inches(0.28), section, 11, True, AMBER)
    add_text(slide, Inches(0.32), Inches(0.34), Inches(10.6), Inches(0.38), title, 22, True, WHITE)
    add_text(slide, Inches(11.4), Inches(0.22), Inches(1.7), Inches(0.36), f"{n} / {SLIDES}", 12, True, SKY, PP_ALIGN.RIGHT)
    add_shape(slide, 0, Inches(7.22), W, Inches(0.28), SLATE)
    add_text(
        slide, Inches(0.32), Inches(7.22), Inches(12.7), Inches(0.28),
        "DT_두산에너빌리티  ·  계획 → 실측 → 근거 있는 변경  ·  ROS 2 Jazzy / Gazebo Harmonic / Unitree Go2",
        10, False, MUTED, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE,
    )


def style_table(table, header_fill=NAVY, sizes=None, aligns=None, bold_col=None):
    sizes = sizes or [11] * len(table.columns)
    aligns = aligns or [PP_ALIGN.LEFT] * len(table.columns)
    for ri, row in enumerate(table.rows):
        for ci, cell in enumerate(row.cells):
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            fill = header_fill if ri == 0 else (ROW_ALT if ri % 2 == 0 else WHITE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            color = WHITE if ri == 0 else TEXT
            is_bold = ri == 0 or (bold_col is not None and ci in bold_col)
            for p in cell.text_frame.paragraphs:
                p.alignment = aligns[ci] if ci < len(aligns) else PP_ALIGN.LEFT
                for run in p.runs:
                    set_run(run, sizes[ci] if ci < len(sizes) else 11, is_bold, color)


def fill_table(table, rows):
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            table.cell(r, c).text = str(val)


def add_tbl(slide, l, t, w, h, data, col_w=None, header_fill=NAVY, sizes=None, aligns=None, bold_col=None):
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, l, t, w, h)
    table = table_shape.table
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            table.columns[i].width = int(w * cw / total)
    fill_table(table, data)
    style_table(table, header_fill, sizes, aligns, bold_col)
    return table


def card(slide, l, t, w, h, fill, title, body, title_color=NAVY, accent=None):
    add_round(slide, l, t, w, h, fill, LINE)
    if accent:
        add_shape(slide, l, t, Inches(0.08), h, accent)
    add_text(slide, l + Inches(0.18), t + Inches(0.08), w - Inches(0.28), Inches(0.32), title, 13, True, title_color)
    add_text(slide, l + Inches(0.18), t + Inches(0.38), w - Inches(0.28), h - Inches(0.46), body, 12, False, TEXT)


# ---------------------------------------------------------------------------
# slides
# ---------------------------------------------------------------------------

def s01(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(s, 0, 0, W, H, NAVY2)
    add_shape(s, 0, 0, Inches(0.16), H, GOLD)
    add_shape(s, 0, Inches(7.18), W, Inches(0.32), NAVY)
    add_text(s, Inches(0.7), Inches(0.38), Inches(12), Inches(0.32),
             "32_DT_두산에너빌리티  ·  PBL ①~⑤ 요약 발표  ·  v3", 13, True, AMBER)
    add_text(s, Inches(0.7), Inches(0.85), Inches(12.2), Inches(0.9),
             "발전소 Digital Twin\n4족 보행 로봇 자율 점검", 36, True, WHITE)
    add_text(s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.36),
             "macOS + Docker arm64  ·  ROS 2 Jazzy  ·  Gazebo Harmonic  ·  Unitree Go2", 16, False, SKY2)

    add_round(s, Inches(0.7), Inches(2.45), Inches(11.9), Inches(1.85), NAVY3)
    add_text(s, Inches(0.95), Inches(2.55), Inches(11.4), Inches(0.28), "프로젝트 한 문장 정의", 12, True, AMBER)
    add_text(
        s, Inches(0.95), Inches(2.88), Inches(11.4), Inches(1.28),
        "우리는 macOS + Docker(arm64 네이티브) 위에서 ROS 2 Jazzy와 Gazebo Harmonic으로 발전소 Digital Twin을 구축하고, "
        "4족 보행 로봇이 험지를 스스로 통과하고(Locomotion), 아날로그 계기와 과열 배관을 눈으로 읽어내며(Inspection AI), "
        "가스 누출원을 추적한 뒤 안전하게 복귀하는(Gas/Safety) 자율 점검 시스템을 개발하여, "
        "그 모든 상태와 결과를 웹 관제 화면에 실시간 시각화한다.",
        14, False, WHITE,
    )

    chips = [
        ("필수 205h", TEAL, 2.05),
        ("보너스 60h", TEAL, 2.15),
        ("Phase 0 실측 후 스택 변경", ORANGE, 3.55),
        ("계획 → 실측 → 근거 있는 변경", GOLD, 3.85),
    ]
    x = Inches(0.7)
    for label, col, ww in chips:
        add_round(s, x, Inches(4.5), Inches(ww), Inches(0.38), col)
        add_text(s, x, Inches(4.5), Inches(ww), Inches(0.38), label, 12, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        x += Inches(ww) + Inches(0.12)

    people = [
        ("도훈", "Locomotion"),
        ("운학", "Inspection AI"),
        ("태우", "순찰 경로 · Critical Path"),
        ("수현", "Virtual Plant · 관제"),
        ("채현", "Gas / Safety"),
    ]
    x = Inches(0.7)
    cw = Inches(2.28)
    for name, role in people:
        add_round(s, x, Inches(5.15), cw, Inches(1.05), NAVY_CARD)
        add_text(s, x + Inches(0.1), Inches(5.22), cw - Inches(0.2), Inches(0.4), name, 16, True, WHITE, PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.1), Inches(5.6), cw - Inches(0.2), Inches(0.48), role, 12, False, SKY, PP_ALIGN.CENTER)
        x += cw + Inches(0.12)

    add_text(s, Inches(0.7), Inches(7.2), Inches(12), Inches(0.26),
             "기준: 과제 요구사항 원문 + PLAN.md (2026-08-31, Phase 0 실측 후 Jazzy/Harmonic 확정)",
             11, False, SKY2, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)


def s02(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "① 과제 분석  ·  1:30", "우리가 해결해야 할 문제는 무엇인가", 2)

    add_round(s, Inches(0.32), Inches(0.95), Inches(12.7), Inches(1.22), LIGHT, LINE)
    add_text(s, Inches(0.5), Inches(1.02), Inches(12.35), Inches(0.28), "핵심 문제", 12, True, NAVY)
    add_text(
        s, Inches(0.5), Inches(1.3), Inches(12.35), Inches(0.78),
        "사람이 들어가기 위험한 구역일수록 점검 주기가 짧고, 계기는 디지털 통신이 안 되는 아날로그 압력계다.\n"
        "→ 작업자 안전 리스크  ·  데이터의 비정형·비실시간성  ·  이상 징후 조기 발견 실패\n"
        "해결: 사람 대신 4족 보행 로봇을 보내고, 결과를 Digital Twin 관제 화면에서 실시간으로 확인한다. (가상 발전소에서 전 과정 검증)",
        13, False, TEXT,
    )

    add_text(s, Inches(0.32), Inches(2.3), Inches(6.4), Inches(0.3), "Stakeholder", 13, True, NAVY)
    add_tbl(
        s, Inches(0.32), Inches(2.62), Inches(6.5), Inches(2.05),
        [
            ["구분", "대상", "기대하는 것"],
            ["1차", "발전소 정비팀", "위험 구역 무인 점검, 계기 수치 자동 기록"],
            ["2차", "로봇솔루션팀", "실기 투입 전 보행·인식 알고리즘 검증 환경"],
            ["3차", "DX추진팀", "설비 상태 데이터의 디지털화 및 관제 통합"],
        ],
        [1.0, 2.0, 3.5],
        sizes=[11, 12, 12],
        aligns=[PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT],
        bold_col=[0],
    )

    add_text(s, Inches(7.05), Inches(2.3), Inches(5.9), Inches(0.3), "최종 결과물", 13, True, NAVY)
    pkgs = [
        ("module1_locomotion", "험지 보행 · MPC · Fall Recovery"),
        ("module2_inspection", "게이지 독해 · 열화상 퓨전"),
        ("module3_gas_safety", "플룸 · Source Seeking · RTH"),
        ("simulation / dashboard", "월드 SDF · 웹 관제"),
        ("docs", "REPORT.md · README.md · LI"),
    ]
    y = Inches(2.62)
    for name, desc in pkgs:
        add_round(s, Inches(7.05), y, Inches(5.95), Inches(0.4), LIGHT, LINE)
        add_text(s, Inches(7.18), y, Inches(2.55), Inches(0.4), name, 11, True, NAVY, valign=MSO_ANCHOR.MIDDLE, name=MONO)
        add_text(s, Inches(9.75), y, Inches(3.1), Inches(0.4), desc, 11, False, TEXT, valign=MSO_ANCHOR.MIDDLE)
        y += Inches(0.44)

    add_round(s, Inches(0.32), Inches(4.85), Inches(12.7), Inches(2.2), BLUE_BG, NAVY)
    add_text(s, Inches(0.5), Inches(4.95), Inches(12.35), Inches(0.3), "모델링 최소 조건  ·  GitHub Repository 1식", 13, True, NAVY)
    mins = [
        ("복도", "45 × 3.5 × 4 m"),
        ("공장", "25 × 18 × 7 m"),
        ("험지 방해물", "복도 3 · 공장 5"),
        ("로봇개", "Unitree Go2 × 1"),
        ("공장기계", "1기 + 가스탱크"),
    ]
    x = Inches(0.5)
    for k, v in mins:
        add_round(s, x, Inches(5.38), Inches(2.4), Inches(0.85), WHITE, LINE)
        add_text(s, x, Inches(5.44), Inches(2.4), Inches(0.32), k, 12, True, MUTED, PP_ALIGN.CENTER)
        add_text(s, x, Inches(5.74), Inches(2.4), Inches(0.4), v, 14, True, NAVY, PP_ALIGN.CENTER)
        x += Inches(2.48)
    add_text(s, Inches(0.5), Inches(6.32), Inches(12.35), Inches(0.5),
             "실제 로봇 하드웨어는 쓰지 않는다. 가상 발전소(Virtual Plant)에서 보행·인식·가스 추적·관제까지 전 과정을 구현하고 정량 지표로 검증한다.",
             13, False, TEXT)


def s03(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "① 과제 분석  ·  1:00  ·  DT-1 직접 답변", "Digital Twin 4요소 ↔ 우리 구현물", 3)
    add_tbl(
        s, Inches(0.28), Inches(0.95), Inches(12.78), Inches(4.55),
        [
            ["DT 구성 요소", "요구사항이 정의한 역할", "우리 구현 (Harmonic)", "실제 발전소 반영"],
            ["Virtual Plant", "모든 데이터가 적용되는\n가상 환경", "simulation/worlds/\n복도·공장 SDF, 계단, 파이프 잔해,\n공장기계, 가스탱크 + heat signature", "공간 치수, 통행 폭,\n계단 경사, 장애물 밀도,\n조명, 배관 표면 온도"],
            ["Robot Twin", "로봇 상태 및\n센서 데이터 생성", "Go2 URDF + gz_ros2_control\ngpu_lidar / rgbd / thermal / IMU\n가스 센서(자체 노드)", "관절 구성, 센서 장착\n위치·시야각, 배터리 소모"],
            ["Inspection Data", "로봇이 수집하는\n설비 상태 데이터", "/inspection/gauge_reading\n/inspection/thermal_alert\n/gas/concentration", "아날로그 계기 수치,\n배관 표면 온도(K),\n가스 농도 분포"],
            ["Monitoring System", "분석 결과 및\n로봇 상태 시각화", "dashboard/\nrosbridge + 웹 대시보드\n(보너스①을 필수에 흡수)", "로봇 위치, 임무 상태,\n점검 이력, 가스 농도맵"],
        ],
        [2.1, 2.5, 4.5, 3.2],
        sizes=[13, 12, 12, 12],
        aligns=[PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT],
        bold_col=[0],
    )
    add_round(s, Inches(0.28), Inches(5.65), Inches(12.78), Inches(1.4), GOLD_BG, GOLD)
    add_text(s, Inches(0.5), Inches(5.75), Inches(12.4), Inches(0.3), "발표에서 한 줄로", 12, True, GOLD)
    add_text(
        s, Inches(0.5), Inches(6.08), Inches(12.4), Inches(0.82),
        "Virtual Plant는 월드와 온도 태그, Robot Twin은 Go2와 센서(+thermal), Inspection Data는 게이지·열화상·가스 토픽,\n"
        "Monitoring System은 rosbridge 웹 대시보드다. 네 요소가 각각 폴더와 토픽에 대응하므로 “역할만 나열한 팀”과 갈린다.",
        14, False, TEXT,
    )


def s04(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "① 과제 분석  ·  2:00  ·  차별화 포인트 1", "Phase 0 실측 → 스택을 바꿨다", 4)

    add_round(s, Inches(0.28), Inches(0.92), Inches(12.78), Inches(0.72), GOLD_BG, GOLD)
    add_text(
        s, Inches(0.48), Inches(0.98), Inches(12.4), Inches(0.58),
        "개발 머신이 macOS(Apple Silicon)라 ROS 2를 네이티브로 돌릴 수 없다.\n"
        "우회가 아니라 부가 목표로 전환: 누구나  docker compose up  한 번으로 재현 가능한 무료 환경.  확정은 계획이 아니라 Phase 0 실측 이후.",
        13, False, TEXT,
    )

    add_tbl(
        s, Inches(0.28), Inches(1.78), Inches(12.78), Inches(2.15),
        [
            ["", "당초 계획", "실측 결과 (2026-08-31)", "최종 확정"],
            ["ROS 배포판", "Humble", "humble-desktop은 amd64 전용\n에뮬레이션 빈 월드 RTF 0.5", "Jazzy"],
            ["시뮬레이터", "Gazebo Classic", "arm64에 Classic도 ros-gz-sim도 없음", "Gazebo Harmonic (gz-sim8)"],
            ["근거", "—", "Jazzy+Harmonic은 arm64에서\nros-gz-sim, gz_ros2_control, Nav2, Foxglove 제공", "arm64 네이티브 실행"],
        ],
        [1.7, 2.2, 5.4, 3.4],
        sizes=[12, 12, 12, 13],
        aligns=[PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.CENTER],
        bold_col=[0, 3],
    )

    card(s, Inches(0.28), Inches(4.1), Inches(6.2), Inches(1.9), WARN_BG, "불리 — 시간 배분을 바꿈",
         "CHAMP 등 Humble 기준 4족 오픈소스는 Jazzy 포팅 검증이 필요하다.\n→ 리스크 1순위. stage1을 1-A / 1-B로 이중화했다. (슬라이드 6)",
         RED, ORANGE)
    card(s, Inches(6.64), Inches(4.1), Inches(6.42), Inches(1.9), TEAL_BG, "유리 — 작업을 줄임",
         "Harmonic에 thermal 카메라가 네이티브다 (type=\"thermal\").\nClassic 기준의 “열화상 합성 플러그인 자작”을 취소했다. 8h → 2h.",
         TEAL, TEAL)

    add_text(s, Inches(0.32), Inches(6.12), Inches(12.7), Inches(0.95),
             "이 의사결정 기록 자체가 ①의 가장 강한 발표 재료다.  “계획대로 실행했다”가 아니라  계획 → 실측 → 근거 있는 변경.\n"
             "전체-1(가장 어려웠던 기술적 문제)의 1순위 후보: Apple Silicon 실행 환경 + Jazzy/Harmonic 전환. 실측 데이터(RTF 0.5)를 보유한다.",
             13, False, TEXT)


def s05(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "② 문제 분해  ·  1:30", "의존 관계 — P2는 모듈이 아니라 인프라", 5)

    nodes = [
        (0.32, 0.95, 12.7, 0.62, NAVY, WHITE, "P0  개발환경   Docker arm64 / ROS 2 Jazzy / Gazebo Harmonic / Foxglove    ← 없으면 아무것도 실행 불가"),
        (0.32, 1.72, 12.7, 0.62, NAVY3, WHITE, "P1  Virtual Plant + Robot Twin    월드 SDF · 온도 태깅 · Go2 URDF · 센서 + ros_gz_bridge"),
        (0.32, 2.49, 12.7, 0.62, ORANGE, WHITE, "P2  Module 1 · Locomotion    Terrain Mapping → stage1 보행 → stage2/3 MPC → Fall Recovery    ★ 55h 병목"),
    ]
    for l, t, w, h, fill, col, txt in nodes:
        add_round(s, Inches(l), Inches(t), Inches(w), Inches(h), fill)
        add_text(s, Inches(l + 0.16), Inches(t), Inches(w - 0.28), Inches(h), txt, 14, True, col, valign=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.32), Inches(3.18), Inches(12.7), Inches(0.28), "여기서부터 병렬", 11, True, MUTED, PP_ALIGN.CENTER)

    tri = [
        (0.32, TEAL, "P3  Module 2  Inspection AI", "게이지 · 열화상-RGB 정합\nViewpoint 순찰 경로"),
        (4.55, GOLD, "P4  Module 3  Gas / Safety", "가우시안 플룸 · Source Seeking\nReturn to Home  (P2에 실질 의존)"),
        (8.78, NAVY, "※ P5는 최종이되 규격은 선행", "토픽/QoS는 P3 착수 시 초안 확정\nP5에서는 검증만"),
    ]
    for l, col, title, body in tri:
        add_round(s, Inches(l), Inches(3.48), Inches(4.08), Inches(1.28), WHITE, col)
        add_shape(s, Inches(l), Inches(3.48), Inches(0.1), Inches(1.28), col)
        add_text(s, Inches(l + 0.22), Inches(3.54), Inches(3.75), Inches(0.36), title, 13, True, col)
        add_text(s, Inches(l + 0.22), Inches(3.9), Inches(3.75), Inches(0.74), body, 12, False, TEXT)

    add_round(s, Inches(0.32), Inches(4.92), Inches(6.2), Inches(0.62), NAVY)
    add_text(s, Inches(0.48), Inches(4.92), Inches(5.9), Inches(0.62),
             "P5  통신 규격 + Monitoring + 통합 시나리오", 14, True, WHITE, valign=MSO_ANCHOR.MIDDLE)
    add_round(s, Inches(6.72), Inches(4.92), Inches(3.0), Inches(0.62), TEAL)
    add_text(s, Inches(6.82), Inches(4.92), Inches(2.8), Inches(0.62), "P6  보너스 ②③④", 13, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_round(s, Inches(9.92), Inches(4.92), Inches(3.08), Inches(0.62), MUTED)
    add_text(s, Inches(10.02), Inches(4.92), Inches(2.88), Inches(0.62), "P7  검증 · 문서화", 13, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

    bullets = [
        "P2가 병목: 걷지 못하면 점검도 가스 추적도 시작 불가. Module 1은 Module 2·3의 전제조건이라 55h.",
        "P1에 ros_gz_bridge가 들어가는 이유: Harmonic은 gz-transport. 브리지를 P2로 미루면 Elevation Map이 통째로 막힌다.",
        "P3/P4 병렬: 가스 확산은 시뮬레이터 무관 자체 노드(가우시안 플룸)라 자원 경합이 없다.",
    ]
    y = Inches(5.68)
    for b in bullets:
        add_text(s, Inches(0.4), y, Inches(12.6), Inches(0.42), "▸  " + b, 13, False, TEXT)
        y += Inches(0.38)


def s06(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "② 문제 분해  ·  1:30  ·  차별화 포인트 2", "Module 1 3단 리스크 분해 + stage1 이중화", 6)

    add_text(s, Inches(0.32), Inches(0.9), Inches(12.7), Inches(0.32),
             "MPC는 실패 확률이 가장 높다.  완주 보장선과 도전선을 분리한다.", 13, True, TEXT)

    add_tbl(
        s, Inches(0.28), Inches(1.24), Inches(12.78), Inches(1.85),
        [
            ["단계", "내용", "성격", "실패 시 영향"],
            ["stage1  champ_tuning", "기성 4족 스택 + 지형 적응 + Trot/Crawl", "완주 보장선 — 필수 요건 충족", "프로젝트 중단"],
            ["stage2  simple_mpc_py", "SRBD 간이 MPC를 Python으로 직접 구현", "이해도 확보", "없음"],
            ["stage3  convex_mpc_cpp", "Convex MPC (MIT Cheetah) C++ · 타임박스 10h", "도전", "없음 · 보고서에 시도·한계"],
        ],
        [2.6, 4.6, 3.4, 2.2],
        sizes=[12, 12, 12, 12],
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER],
        bold_col=[0],
    )

    add_round(s, Inches(0.28), Inches(3.22), Inches(12.78), Inches(0.7), WARN_BG, ORANGE)
    add_text(
        s, Inches(0.48), Inches(3.28), Inches(12.4), Inches(0.58),
        "Jazzy 전환으로 생긴 새 리스크: CHAMP는 Humble 기준이라 Jazzy에서 동작한다는 보장이 없다.\n"
        "완주 보장선이 흔들리면 프로젝트 전체가 흔들린다.  그래서 stage1을 이중화했다.  판정은 P2 착수 전(P0 말~P1 초).",
        13, False, TEXT,
    )

    card(s, Inches(0.28), Inches(4.06), Inches(6.2), Inches(1.55), BLUE_BG, "1-A  우선 경로",
         "CHAMP 커뮤니티 Jazzy 포크를 검증한 뒤 채택.\n실패하면 즉시 1-B로 전환한다.", NAVY, NAVY)
    card(s, Inches(6.64), Inches(4.06), Inches(6.42), Inches(1.55), TEAL_BG, "1-B  대체 경로",
         "gz_ros2_control + 자체 게이트 생성기\n(정적 보행 궤적 + 다리 역기구학)", TEAL, TEAL)

    add_round(s, Inches(0.28), Inches(5.75), Inches(12.78), Inches(1.3), LIGHT, NAVY)
    add_text(s, Inches(0.5), Inches(5.82), Inches(12.4), Inches(0.28), "1-B는 손실이 아니다", 13, True, NAVY)
    add_text(
        s, Inches(0.5), Inches(6.12), Inches(12.4), Inches(0.82),
        "발끝 궤적·역기구학 코드는 stage2 간이 MPC와 그대로 공유된다. 1-B로 가면 stage1과 stage2가 한 코드베이스로 합쳐져 총 작업량이 줄고,\n"
        "“왜 직접 만들었는가”라는 심층 인터뷰 답변이 강해진다. 대신 stage1 튜닝 15h 일부가 구현으로 옮겨가므로 stage3(도전) 포기 가능성을 미리 연다.",
        13, False, TEXT,
    )


def s07(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "③ 학습 주제 도출  ·  1:30", "Learning Issue 29개 — 우선순위가 학습 순서다", 7)

    add_text(s, Inches(0.32), Inches(0.9), Inches(12.7), Inches(0.3),
             "판단 기준  ① 다른 작업의 선행조건인가  ② 정량 목표에 직결되는가  ③ 실패 시 대체 수단이 있는가", 13, True, TEXT)

    prios = [
        (0.28, NAVY, "P0  15개", "없으면 프로젝트가 멈춤\nPhase 0~1에 집중 학습"),
        (4.55, GOLD, "P1", "품질·정량 목표를 좌우\n해당 Phase 착수 직전에"),
        (8.82, MUTED, "P2", "보너스·심화\nPhase 6 직전에"),
    ]
    for l, col, title, body in prios:
        add_round(s, Inches(l), Inches(1.28), Inches(4.08), Inches(1.15), WHITE, col)
        add_shape(s, Inches(l), Inches(1.28), Inches(0.1), Inches(1.15), col)
        add_text(s, Inches(l + 0.22), Inches(1.36), Inches(3.7), Inches(0.36), title, 16, True, col)
        add_text(s, Inches(l + 0.22), Inches(1.72), Inches(3.7), Inches(0.6), body, 13, False, TEXT)

    add_text(s, Inches(0.32), Inches(2.58), Inches(12.7), Inches(0.28), "P0 핵심 (신설 4개가 v2 대비 가장 큰 변화)", 13, True, NAVY)

    lis = [
        ("LI-01", "환경", "Docker arm64 Jazzy 이미지 · Foxglove/noVNC"),
        ("LI-03", "환경/제어", "Humble→Jazzy 이식성 · 포팅 판정 절차  ★신설"),
        ("LI-04", "시스템", "ros_gz_bridge / ros_gz_image  ★신설"),
        ("LI-07", "환경/비전", "Harmonic thermal 센서 · heat signature  ★신설"),
        ("LI-12", "제어", "자체 게이트 생성기 (stage1-B)  ★신설"),
        ("LI-08·09", "제어", "PointCloud 필터 · Elevation Map · 발 디딤점"),
        ("LI-15", "시스템", "Nav2 Jazzy · costmap · Behavior Tree"),
        ("LI-16·18", "비전", "게이지 독해 · 열화상-RGB Extrinsic 정합"),
        ("LI-20·21", "시스템", "Gaussian Plume · Source Seeking / Infotaxis"),
        ("LI-23", "시스템", "TF2 · use_sim_time · message_filters"),
    ]
    y = Inches(2.9)
    x0 = Inches(0.28)
    for i, (code, tr, desc) in enumerate(lis):
        col = i % 2
        x = x0 + Inches(col * 6.4)
        yy = y + Inches((i // 2) * 0.42)
        add_round(s, x, yy, Inches(6.22), Inches(0.38), LIGHT, LINE)
        add_text(s, x + Inches(0.1), yy, Inches(1.15), Inches(0.38), code, 11, True, NAVY, valign=MSO_ANCHOR.MIDDLE, name=MONO)
        add_text(s, x + Inches(1.28), yy, Inches(1.35), Inches(0.38), tr, 11, True, TEAL, valign=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(2.65), yy, Inches(3.45), Inches(0.38), desc, 11, False, TEXT, valign=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.32), Inches(6.95) - Inches(0.88), Inches(12.7), Inches(0.7),
             "삭제: “Gazebo 커스텀 센서 플러그인 작성법” — 네이티브 thermal 사용으로 불필요.\n"
             "P1 예시: LI-13 MPC, LI-14 Fall Recovery, LI-17 저조도, LI-24 rosbridge.  P2: LI-26~29 RL·다중로봇·EC2.",
             13, False, MUTED)


def s08(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "④ 자기주도 학습  ·  1:00", "트랙별 배분 · 담당 산출물 · 학습 정리 템플릿", 8)

    add_tbl(
        s, Inches(0.28), Inches(0.95), Inches(7.55), Inches(2.85),
        [
            ["이름", "담당", "맞닿는 트랙"],
            ["도훈", "Locomotion control", "제어 LI-08~14"],
            ["운학", "Inspection AI", "비전 LI-16~18"],
            ["태우", "순찰 경로 · Critical Path", "비전 LI-19 · /patrol_mission"],
            ["수현", "Gazebo Virtual Plant + 관제", "환경 LI-05~07 · 시스템 LI-24"],
            ["채현", "Gas / Safety", "시스템 LI-20~22"],
        ],
        [1.3, 3.2, 3.0],
        sizes=[12, 12, 12],
        aligns=[PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT],
        bold_col=[0],
    )

    add_tbl(
        s, Inches(8.0), Inches(0.95), Inches(5.05), Inches(2.85),
        [
            ["트랙", "산출물"],
            ["환경", "docker/  simulation/"],
            ["제어", "module1_locomotion/"],
            ["비전", "module2_inspection/"],
            ["시스템", "module3_gas_safety/\ndashboard/  통신 규격"],
            ["공통", "REPORT.md Sim2Real"],
        ],
        [1.4, 3.6],
        sizes=[12, 12],
        aligns=[PP_ALIGN.CENTER, PP_ALIGN.LEFT],
        bold_col=[0],
    )

    add_round(s, Inches(0.28), Inches(3.95), Inches(6.3), Inches(3.05), LIGHT, NAVY)
    add_text(s, Inches(0.48), Inches(4.05), Inches(5.95), Inches(0.32), "docs/learning/LI-XX.md", 14, True, NAVY, name=MONO)
    add_text(
        s, Inches(0.48), Inches(4.42), Inches(5.95), Inches(2.4),
        "1. 핵심 개념 요약 (5줄 이내)\n"
        "2. 우리 코드의 어디에 들어가는가\n    파일 경로 / 입력 토픽 / 출력 토픽\n"
        "3. 선택한 방법과 근거 (대안 대비)\n"
        "4. 예상 어려움 / 미해결 질문\n"
        "5. 출처 (URL · 논문 · 접근일)\n"
        "6. 이 LI가 답하는 동료평가 질문 번호",
        13, False, TEXT,
    )

    add_round(s, Inches(6.75), Inches(3.95), Inches(6.3), Inches(3.05), GOLD_BG, GOLD)
    add_text(s, Inches(6.95), Inches(4.05), Inches(5.95), Inches(0.32), "규칙 두 가지", 14, True, GOLD)
    add_text(
        s, Inches(6.95), Inches(4.45), Inches(5.95), Inches(2.35),
        "① 개념 요약보다 “우리 코드의 어느 파일에\n    어떻게 들어가는가”를 반드시 한 줄 이상.\n\n"
        "② 6번(대응 평가 질문)을 채운다.\n    이 문서가 그대로 심층 인터뷰 답변 초안.\n\n"
        "태우의 Viewpoint/TSP(LI-19)는 운학과\n/patrol_mission Action · 비용함수 입력을 P3 전에 맞춘다.",
        13, False, TEXT,
    )


def s09(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "⑤ 실행계획  ·  2:00", "Phase 로드맵 205h + 보너스 60h  ·  평가 지표", 9)

    add_tbl(
        s, Inches(0.22), Inches(0.92), Inches(8.35), Inches(4.55),
        [
            ["Phase", "시간", "완료 기준", "평가"],
            ["P0 개발환경", "15h", "compose up → 빈 월드, RTF ≥ 0.8", "—"],
            ["P1 Virtual Plant", "25h", "최소 모델링 + ROS에서 전 센서 수신", "DT-1"],
            ["P2 Locomotion", "55h", "험지·계단 통과, 넘어짐 시 자동 기립", "M1-1~3"],
            ["P3 Inspection", "29h", "이동 중 판독 오차 ≤5%, 과열 검출", "M2-1~3"],
            ["P4 Gas/Safety", "30h", "임의 누출원 추적, 비상 시 복귀", "M3-1,2"],
            ["P5 통합+관제", "30h", "스크립트 1개 시나리오 + 대시보드", "통합·DT-2"],
            ["P6 보너스", "60h", "RL · 다중로봇 · 클라우드 관제", "DT-3"],
            ["P7 문서화", "15h", "제출 요건 7종", "전체-1,2"],
            ["예비", "6h", "stage1 이중화 버퍼 (P3 절감분)", "—"],
        ],
        [2.15, 0.85, 3.85, 1.5],
        sizes=[11, 11, 11, 11],
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.CENTER],
        bold_col=[0, 1],
    )

    add_text(s, Inches(8.72), Inches(0.92), Inches(4.35), Inches(0.3), "정량 목표 (P7에서 측정)", 13, True, NAVY)
    metrics = [
        ("빈 월드 RTF", "≥ 0.8"),
        ("통합 시나리오 RTF", "≥ 0.6"),
        ("험지 넘어짐 / 10회", "≤ 1회"),
        ("Fall Recovery", "≥ 90%"),
        ("게이지 판독 오차", "≤ 5%"),
        ("열화상-RGB 정합", "≤ 10 px"),
        ("누출원 탐지 / 거리", "≥80% / ≤1m"),
        ("Return to Home", "100% (20%)"),
        ("대시보드 E2E 지연", "≤ 500 ms"),
    ]
    y = Inches(1.26)
    for k, v in metrics:
        add_round(s, Inches(8.72), y, Inches(4.35), Inches(0.38), LIGHT, LINE)
        add_text(s, Inches(8.84), y, Inches(2.55), Inches(0.38), k, 11, False, TEXT, valign=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(11.4), y, Inches(1.55), Inches(0.38), v, 12, True, NAVY, PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE)
        y += Inches(0.42)

    add_round(s, Inches(0.22), Inches(5.6), Inches(12.9), Inches(1.45), BLUE_BG, NAVY)
    add_text(s, Inches(0.42), Inches(5.68), Inches(12.5), Inches(0.28), "시간 배분의 이유", 13, True, NAVY)
    add_text(
        s, Inches(0.42), Inches(6.0), Inches(12.5), Inches(1.0),
        "P2 55h가 최대: 보행이 나머지 모듈의 전제조건.  P3는 thermal 플러그인 자작 취소로 6h 절감 → 예비 6h로 P2(stage1 이중화)에 이관.\n"
        "P1은 ros_gz_bridge(+3h)와 발열 오브젝트 태깅(+2h)이 늘어 25h.  Phase 완료 기준을 통과하기 전에는 다음 Phase에 착수하지 않는다.\n"
        "데모 영상은 각 Phase 완료 시점에 바로 녹화한다. (P2 험지 / P3 게이지+퓨전 / P4 추적 경로 / P5 통합+대시보드)",
        13, False, TEXT,
    )


def s10(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "⑤ 실행계획  ·  1:00", "리스크 관리 + 보너스 범위 결정", 10)

    add_tbl(
        s, Inches(0.22), Inches(0.92), Inches(12.9), Inches(3.55),
        [
            ["#", "리스크", "대응"],
            ["1", "CHAMP 등 4족 오픈소스 Jazzy 미지원", "포크 검증 → 실패 시 1-B. 궤적을 stage2와 공유. P2 전 판정"],
            ["2", "Apple Silicon Gazebo 성능 저하", "물리 스텝 조정, headless+Foxglove. 최악 시 EC2 CPU"],
            ["3", "ros_gz_bridge 이미지 지연", "ros_gz_image 압축, 필요 토픽만 선별 브리징"],
            ["4", "stage3 C++ MPC 미완성", "영향 없음(stage1이 요건 충족). 타임박스 10h"],
            ["5", "제공 에셋 포맷 비호환", "Blender 변환 우선. 불가 시 자체 모델링(과제 허용)"],
            ["6", "RL 정책의 Gazebo 이식 실패", "도메인 랜덤화. 실패 시 학습 곡선만 보고"],
            ["7", "EC2 비용 초과", "스팟 + 자동 중지 + $30 예산 알람"],
            ["8", "환경 재전환 유혹", "Phase 2 이후 Isaac Sim 전환 금지"],
        ],
        [0.55, 4.4, 7.95],
        sizes=[11, 12, 12],
        aligns=[PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT],
        bold_col=[0],
    )

    add_text(s, Inches(0.28), Inches(4.58), Inches(12.7), Inches(0.28), "보너스 범위", 13, True, NAVY)
    bonuses = [
        (0.22, MUTED, "① 실시간 대시보드", "필수에 흡수 (P5)\nMonitoring이 이미 필수"),
        (3.5, TEAL, "③ RL 보행  수행 30h", "EC2 스팟. Sim2Sim 이식이\nDT-3(Sim2Real)의 실전 답변"),
        (6.78, TEAL, "② 다중 로봇  수행 15h", "네임스페이스 분리만으로\n구현 가능. 투자 대비 효과"),
        (10.06, TEAL, "④ 클라우드 관제  15h", "① 대시보드를 EC2에 배포.\n한계비용이 낮음"),
    ]
    for l, col, title, body in bonuses:
        add_round(s, Inches(l), Inches(4.9), Inches(3.12), Inches(1.15), WHITE, col)
        add_shape(s, Inches(l), Inches(4.9), Inches(0.08), Inches(1.15), col)
        add_text(s, Inches(l + 0.18), Inches(4.96), Inches(2.85), Inches(0.36), title, 12, True, col)
        add_text(s, Inches(l + 0.18), Inches(5.32), Inches(2.85), Inches(0.66), body, 12, False, TEXT)

    add_text(s, Inches(0.28), Inches(6.15), Inches(12.7), Inches(0.9),
             "협업: Git Flow (main / develop / feature/module{n}-{task}) · Issue+Kanban · 주 1회 통합 빌드 + LI 세미나 10분.\n"
             "의사결정 기록: Phase 0의 Jazzy 전환처럼 계획을 바꾼 판단은 docs/decisions/에 실측 근거와 함께 남긴다. REPORT.md 소재로 직결.",
             13, False, TEXT)


def s11(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "부록 A  ·  1:00  ·  통합-1 · 통합-2 · DT-2", "ROS 2 통신 구조 — Harmonic 브리지가 핵심", 11)

    add_round(s, Inches(0.28), Inches(0.92), Inches(12.78), Inches(1.55), NAVY2)
    add_text(s, Inches(0.48), Inches(1.0), Inches(12.4), Inches(0.28), "A-0  브리지 계층 (v3 신설 — Classic에는 없던 단계)", 13, True, AMBER)
    add_text(
        s, Inches(0.48), Inches(1.32), Inches(12.4), Inches(0.95),
        "Gazebo Harmonic 센서는 gz-transport로 발행된다. ROS 2 노드는 직접 구독할 수 없다.\n"
        "[gz-transport]  →  ros_gz_bridge (LiDAR, IMU, JointState, Clock, cmd_vel)\n"
        "                 →  ros_gz_image  (RGB, Depth, Thermal · 압축)   →  [ROS 2 Jazzy] sensor_msgs\n"
        "설계 원칙: 필요 토픽만 선별 브리징.  /clock은 gz→ROS 단방향. 전 노드 use_sim_time: true.",
        13, False, WHITE,
    )

    add_tbl(
        s, Inches(0.22), Inches(2.6), Inches(7.85), Inches(2.55),
        [
            ["방식", "사용 대상", "이유"],
            ["Topic", "센서 스트림, 상태, 점검 결과", "다수 구독 · 연속·비동기"],
            ["Service", "/set_gait_mode, /trigger_gauge_read", "즉시 응답, 짧은 요청-응답"],
            ["Action", "/patrol_mission, /seek_gas_source\n/return_to_home, /fall_recovery", "수 초~수 분, 진행률·취소"],
        ],
        [1.5, 3.5, 2.85],
        sizes=[12, 12, 12],
        aligns=[PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT],
        bold_col=[0],
    )

    add_text(s, Inches(8.25), Inches(2.6), Inches(4.8), Inches(0.28), "QoS 한 줄", 13, True, NAVY)
    qos = [
        ("센서 스트림", "BestEffort  ·  드롭 허용"),
        ("점검 결과 · 상태", "Reliable + TransientLocal"),
        ("/gait_mode, /robot/status", "늦게 붙은 대시보드도 마지막 값"),
        ("열화상 이미지", "L16, 9Hz  ·  raw×resolution=K"),
    ]
    y = Inches(2.92)
    for k, v in qos:
        add_round(s, Inches(8.25), y, Inches(4.8), Inches(0.5), LIGHT, LINE)
        add_text(s, Inches(8.38), y, Inches(4.55), Inches(0.26), k, 11, True, NAVY, valign=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(8.38), y + Inches(0.22), Inches(4.55), Inches(0.24), v, 11, False, TEXT)
        y += Inches(0.54)

    add_round(s, Inches(0.22), Inches(5.28), Inches(12.9), Inches(1.75), LIGHT, NAVY)
    add_text(s, Inches(0.42), Inches(5.36), Inches(12.5), Inches(0.28), "병목·충돌 방지  (통합-1 / 통합-2)", 13, True, NAVY)
    add_text(
        s, Inches(0.42), Inches(5.68), Inches(12.5), Inches(1.22),
        "1. 브리지 선별 — 전 토픽 브리징은 그 자체가 병목. 이미지는 ros_gz_image 압축.\n"
        "2. 대역폭 분리 — 원본 이미지는 로컬 노드만, 대시보드는 CompressedImage + 결과 메시지. rosbridge throttle.\n"
        "3. 실행기 분리 — 제어 루프 200Hz와 비전 추론을 별도 프로세스. 비전은 MultiThreadedExecutor + Reentrant.\n"
        "4. 시간 동기화 — TF2  map→odom→base_link→{lidar, rgb, thermal, gas}. RGB 15Hz vs 열화상 9Hz는 ApproximateTimeSynchronizer.\n"
        "5. 네임스페이스 — 보너스② 대비 처음부터 /robot1/… 구조.  대시보드는 수신 시각이 아니라 header.stamp 기준 정렬.",
        13, False, TEXT,
    )


def s12(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "발표 서사  ·  Q&A 대비", "두 축만 기억하면 된다", 12)

    card(s, Inches(0.28), Inches(0.95), Inches(6.3), Inches(2.55), GOLD_BG, "축 1  ·  슬라이드 4",
         "계획 → 실측 → 근거 있는 변경\nHumble+Classic을 종이 위에서 고수하지 않았다.\nRTF 0.5를 재고, Jazzy+Harmonic으로 바꿨다.\n전체-1의 1순위 후보.", GOLD, GOLD)
    card(s, Inches(6.75), Inches(0.95), Inches(6.3), Inches(2.55), WARN_BG, "축 2  ·  슬라이드 6",
         "완주 보장선과 도전선을 분리했다.\nstage1을 1-A/1-B로 이중화했고,\n1-B 코드는 stage2 MPC와 공유한다.\n실패해도 프로젝트가 멈추지 않는다.", ORANGE, ORANGE)

    add_text(s, Inches(0.32), Inches(3.65), Inches(12.7), Inches(0.3), "동료평가 질문 ↔ 오늘 슬라이드", 13, True, NAVY)
    add_tbl(
        s, Inches(0.22), Inches(3.98), Inches(12.9), Inches(2.05),
        [
            ["질문", "오늘 답하는 곳"],
            ["DT-1  4요소 역할", "슬라이드 3  매핑 표"],
            ["DT-2 · 통합-1·2  전달·동기화·병목", "슬라이드 11  브리지 / QoS / 실행기 분리"],
            ["DT-3  Sim2Real", "보너스③ Sim2Sim (PhysX→DART) + 부록 B"],
            ["전체-1  가장 어려웠던 문제", "슬라이드 4  Apple Silicon + 스택 전환 실측"],
        ],
        [4.6, 8.3],
        sizes=[13, 13],
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT],
        bold_col=[0],
    )

    add_text(
        s, Inches(0.32), Inches(6.15), Inches(12.7), Inches(0.9),
        "데모 영상 확보: P2 종료 시 험지 통과  ·  P3 종료 시 게이지 오버레이 + 열화상 퓨전  ·  P4 종료 시 추적 경로  ·  P5 종료 시 통합+대시보드.\n"
        "발표 직전에 몰리지 않게, 각 Phase 완료 시점에 바로 녹화한다.",
        13, False, TEXT,
    )


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    s01(prs)
    s02(prs)
    s03(prs)
    s04(prs)
    s05(prs)
    s06(prs)
    s07(prs)
    s08(prs)
    s09(prs)
    s10(prs)
    s11(prs)
    s12(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
