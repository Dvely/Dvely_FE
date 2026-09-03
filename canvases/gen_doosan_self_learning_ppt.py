#!/usr/bin/env python3
"""④ 자기주도 학습 — 팀원별 LI 조사 핵심·출처 PPT."""

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

OUT = "/Users/kimtaewoo/Documents/GitHub/Dvely_FE/canvases/doosan-pbl-self-learning.pptx"
FONT = "Apple SD Gothic Neo"
MONO = "Menlo"


def C(h):
    return RGBColor((h >> 16) & 0xFF, (h >> 8) & 0xFF, h & 0xFF)


NAVY = RGBColor(0x1E, 0x4A, 0x7A)
GOLD = C(0xB45309)
TEXT = C(0x1A1A1A)
MUTED = C(0x4B5563)
LINE = C(0xD1D5DB)
WHITE = C(0xFFFFFF)
LIGHT = C(0xF4F7FB)
GOLD_BG = C(0xFFFBEB)
AMBER = C(0xFCD34D)
SKY = C(0xBFDBFE)
SLATE = C(0xF1F5F9)
ROW_ALT = C(0xF8FAFC)
BLUE_BG = C(0xE8F0F8)

W = Inches(13.333)
H = Inches(7.5)
SLIDES = 7


def ea_font(run, name=FONT):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", name)


def set_run(run, size, bold=False, color=TEXT, name=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    ea_font(run, name)


def add_shape(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def add_round(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.adjustments[0] = 0.08
    sh.shadow.inherit = False
    return sh


def tb(shape, text, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE, name=FONT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.anchor = valign
    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.clear()
        run = p.add_run()
        run.text = line
        set_run(run, size, bold, color, name)
    return tf


def add_text(slide, l, t, w, h, text, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, name=FONT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tb(box, text, size, bold, color, align, valign, name)
    return box


def header(slide, section, title, n):
    add_shape(slide, 0, 0, W, Inches(0.78), NAVY)
    add_shape(slide, 0, 0, Inches(0.12), Inches(0.78), GOLD)
    add_text(slide, Inches(0.32), Inches(0.08), Inches(10.4), Inches(0.28), section, 11, True, AMBER)
    add_text(slide, Inches(0.32), Inches(0.34), Inches(10.6), Inches(0.38), title, 20, True, WHITE)
    add_text(slide, Inches(11.4), Inches(0.22), Inches(1.7), Inches(0.36), f"{n} / {SLIDES}", 12, True, SKY, PP_ALIGN.RIGHT)
    add_shape(slide, 0, Inches(7.22), W, Inches(0.28), SLATE)
    add_text(
        slide, Inches(0.32), Inches(7.22), Inches(12.7), Inches(0.28),
        "④ 자기주도 학습  ·  docs/learning/LI-XX.md  ·  접근일 2026-09-03  ·  Jazzy / Harmonic / Go2",
        10, False, MUTED, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE,
    )


def style_table(table, header_fill=NAVY, sizes=None, aligns=None, bold_col=None):
    sizes = sizes or [11] * len(table.columns)
    aligns = aligns or [PP_ALIGN.LEFT] * len(table.columns)
    for ri, row in enumerate(table.rows):
        for ci, cell in enumerate(row.cells):
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            fill = header_fill if ri == 0 else (ROW_ALT if ri % 2 == 0 else WHITE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            color = WHITE if ri == 0 else TEXT
            is_bold = ri == 0 or (bold_col is not None and ci in bold_col)
            for p in cell.text_frame.paragraphs:
                p.alignment = aligns[ci] if ci < len(aligns) else PP_ALIGN.LEFT
                for run in p.runs:
                    set_run(run, sizes[ci] if ci < len(sizes) else 11, is_bold, color)


def fill_table(table, rows):
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            table.cell(r, c).text = str(val)


def add_tbl(slide, l, t, w, h, data, col_w=None, sizes=None, aligns=None, bold_col=None):
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, l, t, w, h)
    table = table_shape.table
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            table.columns[i].width = int(w * cw / total)
    fill_table(table, data)
    style_table(table, NAVY, sizes, aligns, bold_col)
    return table


def li_card(slide, l, t, w, h, li, title, core, code, src, eval_q):
    add_round(slide, l, t, w, h, LIGHT, LINE)
    add_shape(slide, l, t, Inches(0.1), h, NAVY)
    add_text(slide, l + Inches(0.22), t + Inches(0.08), w - Inches(0.36), Inches(0.32), f"{li}  {title}", 14, True, NAVY)
    add_text(slide, l + Inches(0.22), t + Inches(0.40), w - Inches(0.36), Inches(0.58), core, 12, False, TEXT)
    add_text(slide, l + Inches(0.22), t + Inches(0.98), w - Inches(0.36), Inches(0.28), code, 11, False, MUTED, name=MONO)
    add_text(slide, l + Inches(0.22), t + Inches(1.26), w - Inches(0.36), Inches(0.42), src, 11, False, TEXT)
    add_text(slide, l + Inches(0.22), t + Inches(1.66), w - Inches(0.36), Inches(0.28), eval_q, 11, True, GOLD)


def person_slide(prs, n, kicker, title, items):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, kicker, title, n)
    gap = Inches(0.12)
    card_h = Inches(1.98)
    y0 = Inches(0.92)
    for i, item in enumerate(items):
        li_card(s, Inches(0.28), y0 + i * (card_h + gap), Inches(12.78), card_h, *item)


def s01(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "④ 자기주도 학습", "팀원별 Learning Issue 3개 — 조사·학습 배분", 1)
    add_tbl(
        s, Inches(0.28), Inches(0.95), Inches(12.78), Inches(3.55),
        [
            ["이름", "역할", "선택한 LI", "산출물"],
            ["도훈", "Locomotion", "LI-09 Elevation Map  ·  LI-10 Trot/Crawl  ·  LI-12 게이트", "module1_locomotion/"],
            ["운학", "Inspection AI", "LI-16 게이지  ·  LI-17 전처리  ·  LI-18 열화상-RGB", "module2_inspection/"],
            ["태우", "순찰 경로 · Critical Path", "LI-19 Viewpoint  ·  TSP  ·  A*  (한 LI를 평가 질문 3개로 분해)", "module2_inspection/\n/patrol_mission"],
            ["수현", "Virtual Plant + 관제", "LI-05 SDF  ·  LI-06 URDF·control  ·  LI-24 대시보드", "docker/  simulation/\ndashboard/"],
            ["채현", "Gas / Safety", "LI-20 플룸  ·  LI-21 Source Seeking  ·  LI-22 RTH", "module3_gas_safety/"],
        ],
        [1.2, 2.4, 5.6, 3.5],
        sizes=[12, 13, 12, 12],
        aligns=[PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT],
        bold_col=[0],
    )
    add_round(s, Inches(0.28), Inches(4.68), Inches(12.78), Inches(2.32), GOLD_BG, GOLD)
    add_text(s, Inches(0.48), Inches(4.80), Inches(12.4), Inches(0.32), "규칙 · 이 장이 답하는 것", 14, True, GOLD)
    add_text(
        s, Inches(0.48), Inches(5.16), Inches(12.4), Inches(1.70),
        "개념만 쓰지 않는다. 각 LI는 파일·입력 토픽·출력 토픽을 한 줄 이상 적는다. 6번(동료평가·심층 질문)을 채우면 이 슬라이드가 인터뷰 초안이다.\n"
        "태우는 표의 LI 번호가 하나(LI-19)다. 동료평가가 Viewpoint 비용 / TSP / A*를 따로 묻기 때문에 조사 노트 3장으로 쪼갠다.\n"
        "P3 착수 전에 태우·운학이 /patrol_mission Action 필드를 맞춘다. MPC·Fall Recovery·Nav2 본구현은 이번 3개 다음 Phase에서 붙인다.",
        14, False, TEXT,
    )


def s02(prs):
    person_slide(
        prs, 2,
        "④ 도훈  ·  Locomotion",
        "LI-09 · LI-10 · LI-12  ·  module1_locomotion/",
        [
            (
                "LI-09",
                "Elevation Map과 발 디딤점",
                "라이다·Depth를 2.5D 높이맵으로 누적한다. 분산이 큰 칸은 디딤 후보에서 빼고, 다리 워크스페이스 안에서 지지다각형이 넓은 칸을 고른다.",
                "코드  elevation/   입력 /lidar/points  /tf   출력 /elevation_map  /footholds",
                "출처  Fankhauser et al., Robot-Centric Elevation Mapping, ICAR 2014  ·  ANYbotics/grid_map  ·  elevation_mapping CPU",
                "평가  Module 1-① 노이즈를 걸러 맵을 만들고 발 디딤을 정하는가  /  심층: 센서 노이즈",
            ),
            (
                "LI-10",
                "Trot / Crawl 안정성 차이와 전환",
                "Trot는 대각선 두 발이 들어 빠르지만 험지에서 넘어진다. Crawl는 세 발 접지로 느리고 안정적이다. 높이맵 거칠기·IMU 피치/롤로 전환한다.",
                "코드  gait/   입력 /elevation_map  /imu   출력 /gait_mode  (trot|crawl)",
                "출처  Raibert, Legged Robots That Balance  ·  ANYmal 보행 모드  ·  Unitree Go2 gait 문서",
                "평가  심층 Module 1: Trot·Crawl 전환 기준 및 방식",
            ),
            (
                "LI-12",
                "자체 게이트 생성기 (stage1-B)",
                "발끝 궤적(사이클로이드/Bezier) + 다리 IK로 관절각을 푼다. CHAMP가 Jazzy·arm64에서 실패하면 이 노드가 보행을 담당한다. MPC는 그 위에 얹는 P1이다.",
                "코드  gait_generator/   입력 /footholds  /gait_mode  /joint_states   출력 /joint_trajectory",
                "출처  chvmp/champ  ·  gz_ros2_control  ·  Unitree Go2 URDF  ·  Di Carlo, MIT Cheetah 3 MPC, ICRA 2018",
                "평가  “CHAMP가 안 되면 무엇으로 걷히는가” 선행 답  /  심층: 가장 어려웠던 이식 문제",
            ),
        ],
    )


def s03(prs):
    person_slide(
        prs, 3,
        "④ 운학  ·  Inspection AI",
        "LI-16 · LI-17 · LI-18  ·  module2_inspection/",
        [
            (
                "LI-16",
                "아날로그 게이지 독해 (오차 ≤5%)",
                "원 검출 → 바늘 직선 → 영점 대비 각도 → 눈금 범위로 환산. OCR은 보조. Hough를 1차로 쓰는 이유는 인터뷰에서 수식으로 답할 수 있어서다. CNN 키포인트는 5%를 못 지킬 때의 2차.",
                "코드  gauge/   입력 /camera/rgb/image   출력 /gauge/reading  /gauge/overlay",
                "출처  OpenCV HoughCircles · HoughLinesP  (docs.opencv.org)  ·  과제 #5 Module 2-1",
                "평가  Module 2-① 바늘·숫자 구분  /  심층: Hough vs Deep Learning 선택 이유",
            ),
            (
                "LI-17",
                "저조도·모션블러 전처리",
                "CLAHE로 국소 대비를 올리고, 보행 진동으로 번진 바늘은 deblur 또는 프레임 스태킹으로 에지를 살린다. LI-16은 전처리 토픽만 구독한다.",
                "코드  preprocess/   입력 /camera/rgb/image   출력 /camera/rgb/image_pp",
                "출처  Pizer et al., CLAHE  ·  OpenCV CLAHE / fastNlMeansDenoising  ·  과제 #5 흔들림·저조도",
                "평가  Module 2-①의 전처리 절",
            ),
            (
                "LI-18",
                "열화상-RGB 정합 (Intrinsic / Extrinsic)",
                "두 카메라는 위치·화각이 다르다. Zhang 캘리브레이션 + Extrinsic. 평면 배관 ROI만 Homography. Harmonic thermal을 별도 링크에 붙이므로 오차가 시뮬에도 실재한다. 합성 플러그인 자작은 하지 않는다.",
                "코드  fusion/   입력 rgb_pp  /camera/thermal/image  /tf   출력 /inspection/thermal_overlay  /thermal_alert",
                "출처  Zhang, Camera Calibration, TPAMI 2000  ·  OpenCV stereoCalibrate  ·  gz-sensors ThermalCamera  ·  tf2",
                "평가  Module 2-② 위치 오차의 수학적 보정  /  심층: 정합 오차 최소화",
            ),
        ],
    )


def s04(prs):
    person_slide(
        prs, 4,
        "④ 태우  ·  순찰 경로 · Critical Path",
        "LI-19를 Viewpoint / TSP / A* 세 장으로 분해  ·  /patrol_mission",
        [
            (
                "LI-19-A",
                "Viewpoint와 인식 비용함수",
                "Viewpoint는 게이지를 바라보는 포즈다. 비용은 법선-시선 각, 거리, 가림, 예상 블러의 역수. 각 게이지 앞 원호에서 후보를 샘플링해 최소 비용 포즈를 고른다. 전공간 최적화는 하지 않는다.",
                "코드  viewpoint/   입력 /inspect_targets · 맵 · 카메라 파라미터   출력 /viewpoints",
                "출처  Next-Best-View / inspection viewpoint 서베이  ·  OpenCV projectPoints  ·  과제 #5 Module 2-3",
                "평가  Module 2-③ 인식이 잘 되는 각도·거리의 가중치  /  심층: Viewpoint 제약",
            ),
            (
                "LI-19-B",
                "방문 순서 — TSP",
                "Viewpoint N개의 순서는 TSP다. 거리 행렬은 유클리드가 아니라 A* 충돌 회피 경로 길이다. N이 십여 개면 2-opt로 충분하다. OR-Tools는 의존성 대비 이득이 적다.",
                "코드  planner/tsp.py   입력 /viewpoints · 거리 행렬   출력 /patrol_order   Action /patrol_mission",
                "출처  Croes, 2-opt, 1958  ·  심층 인터뷰 Module 2: TSP를 무엇으로 썼는가",
                "평가  심층 Module 2 순찰 경로 최적화 알고리즘",
            ),
            (
                "LI-19-C",
                "Viewpoint 사이 연결 — A*",
                "격자·costmap에서 A*로 정적 장애물을 피해 연결한다. 동적 장애물·RTH는 채현 Nav2. 태우는 정적 공장·복도 맵을 전제한다. P3 전에 운학과 /patrol_mission 필드를 맞춘다.",
                "코드  planner/astar.py 또는 Nav2 planner   출력 /patrol_path  (nav_msgs/Path)",
                "출처  Hart, Nilsson, Raphael, A*, 1968  ·  Nav2 Jazzy planner (docs.nav2.org)",
                "평가  Module 2-③ 경로 절  ·  운학과 인터페이스 합의",
            ),
        ],
    )


def s05(prs):
    person_slide(
        prs, 5,
        "④ 수현  ·  Virtual Plant + 관제",
        "LI-05 · LI-06 · LI-24  ·  docker/ · simulation/ · dashboard/",
        [
            (
                "LI-05",
                "Harmonic SDF 월드와 제공 에셋",
                "SDF는 공장·조명·가스 영역, URDF는 로봇. 최소: 복도 45×3.5×4 m, 공장 25×18×7 m, 방해물 복도 3·공장 5. 제공 zip의 DAE/STL은 Blender에서 스케일·충돌을 고친다. Isaac Sim으로 본 시뮬을 바꾸지 않는다.",
                "코드  simulation/worlds/plant.sdf  ·  simulation/models/",
                "출처  gazebosim.org/docs/harmonic  ·  PBL_AssetPackge.zip  ·  과제 #2-7 모델링 최소 조건",
                "평가  DT-① Virtual Plant가 실제 발전소의 무엇을 반영하는가",
            ),
            (
                "LI-06",
                "Go2 URDF와 gz_ros2_control",
                "Robot Twin은 URDF + 센서 프레임이다. optical frame(_optical, z-forward)을 여기서 고정하지 않으면 운학 정합이 전부 틀린다. RGB·thermal·LiDAR·가스를 링크에 부착. thermal은 Harmonic 네이티브.",
                "코드  simulation/models/go2/   spawn: ros_gz_sim  ·  관절 명령은 도훈",
                "출처  ros-controls/gz_ros2_control  ·  gazebosim/ros_gz 메시지 대응표  ·  Unitree Go2 description",
                "평가  DT-② Virtual Plant ↔ Robot Twin 데이터 흐름",
            ),
            (
                "LI-24",
                "rosbridge 웹 대시보드 (보너스①을 P5에 흡수)",
                "브라우저가 WebSocket으로 토픽을 구독한다. 원본 이미지는 로컬, 대시보드는 CompressedImage + 결과 메시지. 로봇 위치·시야·게이지·가스 농도. E2E ≤500 ms는 P5-4에서 측정.",
                "코드  dashboard/   구독 pose · /gauge/reading · /thermal_alert · /gas/concentration · /battery   포트 9090",
                "출처  RobotWebTools/rosbridge_suite  ·  Foxglove (개발 GUI 우회)  ·  과제 #1-2 · #6 관제 화면",
                "평가  DT-② 관제 시각화  /  통합-② 지연·실시간성",
            ),
        ],
    )


def s06(prs):
    person_slide(
        prs, 6,
        "④ 채현  ·  Gas / Safety",
        "LI-20 · LI-21 · LI-22  ·  module3_gas_safety/",
        [
            (
                "LI-20",
                "Gaussian Plume 확산",
                "풍하 거리에서 농도가 횡·연직 가우시안으로 퍼진다. σ는 Pasquill-Gifford로 커지고, 난류가 만든 구멍이 LI-21의 국소 최적이 된다. 닫힌 공장은 벽 마스크. CFD는 범위 밖. Gazebo 입자 대신 자체 노드가 농도를 샘플링한다.",
                "코드  plume/   파라미터: 누출원 · Q · 바람   출력 /gas/concentration",
                "출처  표준 Gaussian Plume (EPA ISC / Sutton)  ·  Pasquill-Gifford σy, σz  ·  과제 #5 Module 3-1",
                "평가  심층 Module 3: 어떤 확산 모델, 실제와의 차이",
            ),
            (
                "LI-21",
                "Source Seeking과 국소 최적 탈출",
                "농도 최대화가 목표이므로 구현은 경사 상승이다(원문 ‘하강’은 보고서에서 바로잡음). 구멍에 갇히면 나선 탐색. Infotaxis는 2차 인용. 1차는 유한차분 구배 + 나선.",
                "코드  seek/   입력 /gas/concentration  /odom   출력 /cmd_vel 또는 Nav2 목표",
                "출처  Vergassola et al., Infotaxis, Nature 2007  ·  Olfactory navigation 서베이  ·  과제 #5 Module 3-2",
                "평가  Module 3-① 구배 → 방향 벡터  /  심층 Local Optima",
            ),
            (
                "LI-22",
                "상태 머신과 Return to Home",
                "Patrol / Seek / Return / Charge. 트리거: 배터리 < 20%, 통신 음영. Return이 Seek·순찰보다 우선. 정적 경로는 태우, 실행·동적 장애물은 Nav2 컨트롤러(LI-15를 여기서 소비).",
                "코드  fsm/   입력 /battery  /comm_health   출력 NavigateToPose  ·  BT에서 Return 우선",
                "출처  Nav2 Behavior Tree · NavigateToPose (docs.nav2.org, Jazzy)  ·  과제 #5 Module 3-3",
                "평가  Module 3-② 복귀 트리거와 경로  /  심층 동적 장애물",
            ),
        ],
    )


def s07(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "④ 출처 북마크", "공통 자료 · 이번 3개에서 뺀 LI · 다음 액션", 7)
    add_tbl(
        s, Inches(0.28), Inches(0.92), Inches(7.7), Inches(4.55),
        [
            ["구분", "자료 (접근일 2026-09-03)"],
            ["공식 문서", "ROS 2 Jazzy docs.ros.org  ·  Gazebo Harmonic  ·  gz-sensors ThermalCamera  ·  Nav2  ·  OpenCV  ·  rosbridge"],
            ["브리지", "gazebosim/ros_gz 메시지 타입 대응표"],
            ["오픈소스", "chvmp/champ  ·  ANYbotics/grid_map  ·  elevation_mapping  ·  gz_ros2_control  ·  Go2 description"],
            ["논문", "Fankhauser 2014  ·  Di Carlo Cheetah 3 2018  ·  Zhang 2000  ·  Hart A* 1968  ·  Vergassola Infotaxis 2007"],
            ["도구", "Foxglove (GUI 우회). Isaac Lab/Genesis는 보너스만. 본 시뮬 Harmonic 유지"],
            ["에셋", "https://download.codysseycampus.kr/PBL_AssetPackge.zip"],
        ],
        [1.6, 6.1],
        sizes=[12, 12],
        aligns=[PP_ALIGN.CENTER, PP_ALIGN.LEFT],
        bold_col=[0],
    )
    add_round(s, Inches(8.16), Inches(0.92), Inches(4.9), Inches(4.55), BLUE_BG, NAVY)
    add_text(s, Inches(8.36), Inches(1.04), Inches(4.5), Inches(0.32), "이번 3개에서 뺀 것", 14, True, NAVY)
    add_text(
        s, Inches(0.48 + 7.88), Inches(1.42), Inches(4.5), Inches(3.85),
        "LI-01·03  도커·이식 — Phase 0 공통\n"
        "LI-02·04·23  통신 규격 — P5 공유\n"
        "LI-07  thermal 태깅 — 수현 SDF + 운학 정합\n"
        "LI-08·11·13·14  필터·CHAMP·MPC·기립\n"
        "            → 도훈, 맵·게이트 다음\n"
        "LI-15  Nav2 본구현 — 채현 RTH가 실사용\n"
        "LI-25~29  Sim2Real·RL·다중로봇·EC2\n"
        "            → 보너스·REPORT.md",
        13, False, TEXT,
    )
    add_text(
        s, Inches(0.28), Inches(5.62), Inches(12.78), Inches(1.4),
        "다음 액션  각자 자기 장(2~6)을 docs/learning/LI-XX.md로 복사하고, 코드 경로를 실제 파일명으로 고친다.\n"
        "태우·운학은 P3 전에 /patrol_mission 필드 한 장으로 합의한다.",
        14, False, TEXT,
    )


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    s01(prs)
    s02(prs)
    s03(prs)
    s04(prs)
    s05(prs)
    s06(prs)
    s07(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
